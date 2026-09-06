#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Дека двумя файлами: PDF для судей и PPTX для проектора.

    npm run pitch:pdf     только PDF
    npm run pitch:pptx    только PPTX
    npm run pitch:deck    оба сразу — так быстрее, слайды снимаются один раз

Оба файла собираются из одних и тех же снимков landing/pitch.html. Это и
есть главное решение этого скрипта, и оно стоило одной переделки.

Как было. PDF печатался браузером напрямую (--print-to-pdf), PPTX
собирался из снимков. Файлы расходились: печать раскладывает страницу по
СВОЕЙ ширине, а не по ширине окна, поэтому в PDF сетки схлопывались в две
колонки, подгон масштаба промахивался мимо настоящей высоты, и текст
наезжал на текст. Проверено 06.09.2026 постранично: на слайде
юнит-экономики верхняя метка и нижняя строка таблицы уходили за край
листа. Глеб сказал прямо: «pdf кривые, херово видно».

Как стало. Снимок — единственный способ увидеть слайд до сборки, и он же
единственный источник для обоих файлов. Что снято, то и в PDF, и в PPTX:
разойтись нечему.

Чем платим. Текст в PDF растровый — не выделяется и не ищется. Для
раздатки на защите это приемлемо: её открывают и смотрят. Взамен получаем
то, чего векторная печать не давала ни разу за три подхода, — гарантию,
что страница выглядит ровно так, как слайд на экране.

Числа деки по-прежнему проверяет npm run check:pitch по разметке
landing/pitch.html, а не по этим файлам.
"""

import hashlib
import html as html_mod
import io
import json
import re
import subprocess
import sys
import tempfile
import time
from http.server import HTTPServer, SimpleHTTPRequestHandler
from pathlib import Path
from threading import Thread

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")

ROOT = Path(__file__).resolve().parent.parent
LANDING = ROOT / "landing"
DECK = LANDING / "pitch.html"
OUT_PDF = LANDING / "RentHUB-pitch.pdf"
OUT_PPTX = LANDING / "RentHUB-pitch.pptx"
PORT = 8901

# Кадр 16:9. Снимаем вдвое крупнее слайда (1920×1080 против 960×540
# читаемых), чтобы на проекторе и при печати на бумагу не рассыпалось.
SHOT_W, SHOT_H = 1920, 1080

# Страница PDF — тот же кадр в точках PostScript: 13.333 × 7.5 дюйма,
# стандартный слайд 16:9. Совпадает с размером слайда PPTX, то есть оба
# файла показывают судьям буквально одно и то же.
PAGE_W, PAGE_H = 960.0, 540.0

BROWSERS = [
    r"C:/Program Files (x86)/Microsoft/Edge/Application/msedge.exe",
    r"C:/Program Files/Microsoft/Edge/Application/msedge.exe",
    r"C:/Program Files/Google/Chrome/Application/chrome.exe",
    r"C:/Program Files (x86)/Google/Chrome/Application/chrome.exe",
    "/usr/bin/google-chrome",
    "/usr/bin/chromium",
]


def find_browser():
    for path in BROWSERS:
        if Path(path).exists():
            return path
    return None


def serve():
    """Локальный сервер вместо file://.

    Из file:// Chromium не берёт часть шрифтов и SVG — слайд снимается в
    системном шрифте, и заметно это только в готовом файле, то есть
    поздно.
    """

    class Handler(SimpleHTTPRequestHandler):
        def __init__(self, *args, **kwargs):
            super().__init__(*args, directory=str(LANDING), **kwargs)

        def log_message(self, *args):
            pass

    server = HTTPServer(("127.0.0.1", PORT), Handler)
    Thread(target=server.serve_forever, daemon=True).start()
    return server


def deck_fingerprint():
    """Отпечаток деки — тот же, что считает scripts/deck.mjs.

    Переводы строк нормализуются: локально файл лежит с CRLF, в
    репозитории с LF, и хеш сырого файла расходился на каждом прогоне CI.
    Проверка, которая падает всегда, учит не читать её вывод.
    """
    text = DECK.read_text(encoding="utf-8").replace("\r\n", "\n")
    return hashlib.sha256(text.encode("utf-8")).hexdigest()


def slide_titles():
    """Названия слайдов — для оглавления PDF.

    Берутся из самой деки: метка слайда («5 Рынок») и его заголовок. Второй
    копии названий в проекте нет, и это то же правило, что с числами —
    расходиться нечему.

    Оглавление в раздатке не украшение. Двадцать одна страница снимками
    листается только подряд: искать в ней слайд про юнит-экономику — значит
    прокручивать всё. С закладками судья открывает нужный за один клик, а
    во время вопросов это и решает, покажем мы ответ или будем искать.
    """
    html = DECK.read_text(encoding="utf-8")
    out = []

    for num, label, span, head in re.findall(
        r'<div class="slide-num"><b>(\d+)</b>\s*([^<]*?)\s*<span>(.*?)</span></div>\s*<h[12]>(.*?)</h[12]>',
        html,
        re.S,
    ):
        head = " ".join(re.sub(r"<[^>]+>", " ", head).split())
        if len(head) > 46:
            head = head[:45].rstrip(" ,—-") + "…"

        # Слайды-продолжения уходят вторым уровнем под свой слайд.
        #
        # Семь из двадцати одного — это вторые и третьи экраны одной мысли:
        # «Рынок» и схема денег, юнит-экономика и бизнес-модель. Плоский
        # список из двадцати одной строки, где «5. Рынок» встречается
        # дважды, читается как ошибка нумерации. Вложенный показывает
        # структуру доклада — и сворачивается до десяти пунктов.
        if "продолжение" in span:
            out.append((2, head))
        else:
            out.append((1, f"{num}. {label.strip()} — {head}"))

    return out


def slide_anchors():
    """id секций по порядку: «#done» → номер страницы в раздатке.

    Внутренние ссылки деки («Что уже работает») в PDF не должны никуда
    вести из интернета — им место внутри файла. Здесь и берётся
    соответствие, по которому сборщик кладёт переход на нужную страницу.
    """
    html = DECK.read_text(encoding="utf-8")
    pages = {}

    for i, match in enumerate(re.finditer(r"<section(?=[ >])([^>]*)>", html)):
        found = re.search(r'id="([^"]+)"', match.group(1))
        if found:
            pages[found.group(1)] = i + 1

    return pages


def shoot(browser, tmp, n):
    """Снимок слайда и координаты его ссылок.

    Оба выхода за один запуск браузера: --screenshot пишет картинку,
    --dump-dom отдаёт готовую разметку в stdout, а в ней — атрибут
    data-links, который страница заполняет после подгона масштаба.
    Второй запуск ради координат удвоил бы время сборки на ровном месте.
    """
    shot = Path(tmp) / f"slide-{n:02d}.png"

    done = subprocess.run(
        [
            browser,
            "--headless=new",
            "--disable-gpu",
            # Свой профиль обязателен: без него запуск не создаёт процесс,
            # а передаёт команду уже открытому браузеру — и снимок не
            # появляется вовсе. На машине, где Edge открыт всегда, это
            # выглядит как «скрипт сломался», хотя сломан не он.
            f"--user-data-dir={Path(tempfile.gettempdir()) / 'renthub-deck-shots'}",
            "--hide-scrollbars",
            f"--window-size={SHOT_W},{SHOT_H}",
            f"--screenshot={shot}",
            # Разметка нужна ради координат ссылок — см. read_links ниже.
            "--dump-dom",
            # Странице нужно досчитать вёрстку и подгон масштаба. Меньше
            # четырёх секунд — и в кадр попадает слайд до подгона.
            "--virtual-time-budget=6000",
            # Отпечаток деки в адресе — не украшение, а лекарство.
            #
            # Профиль браузера переиспользуется между запусками (иначе
            # Chromium каждый раз заводит новый и тратит на это секунды), а
            # вместе с профилем переиспользуется КЕШ. 06.09.2026 это поймано
            # с поличным: дека была правлена, сборка отработала, файлы
            # обновились — и в них попали слайды предыдущей версии. Снаружи
            # не отличить: страниц столько же, отпечаток рядом с файлом
            # свежий (он считается по исходнику, а не по снятому).
            #
            # Адрес со свежим отпечатком кеш не находит и идёт на сервер. Не
            # менялась дека — не меняется адрес, и кеш работает как раньше.
            f"http://127.0.0.1:{PORT}/pitch.html?slide={n}&v={deck_fingerprint()[:12]}",
        ],
        capture_output=True,
        timeout=90,
    )

    # Chromium иногда возвращает управление раньше, чем допишет файл.
    for _ in range(20):
        if shot.exists() and shot.stat().st_size > 0:
            break
        time.sleep(0.3)

    if not shot.exists():
        return None, [], [], None

    links, lines, viewport = read_links(done.stdout.decode("utf-8", "replace"))
    return shot, links, lines, viewport


def read_links(dom):
    """Прямоугольники ссылок слайда из разметки, снятой браузером.

    Пусто — не ошибка: на большинстве слайдов ссылок нет. Ошибкой было бы
    молча потерять их там, где они есть, поэтому разбор ничего не угадывает:
    нет атрибута или он не разбирается — возвращаем пусто и идём дальше, а
    слайд остаётся картинкой, как и был.
    """
    def attr(name):
        found = re.search(r'data-' + name + r'="([^"]*)"', dom)
        if not found:
            return []
        try:
            return json.loads(html_mod.unescape(found.group(1)))
        except (ValueError, TypeError):
            return []

    links = attr("links")
    lines = attr("text")

    # Вьюпорт страницы. Он МЕНЬШЕ кадра снимка: рамка окна съедает по
    # 24 пикселя вширь и 92 ввысь, а --screenshot отдаёт окно целиком.
    # Слайд центрирован, поэтому разница делится поровну — по ней и
    # правятся координаты ссылок.
    viewport = None
    vp_match = re.search(r'data-vp="(\d+)x(\d+)', dom)
    if vp_match:
        viewport = (int(vp_match.group(1)), int(vp_match.group(2)))

    return links, lines, viewport


# Ширина картинки внутри PDF. Снимок делается в 1920 — это ширина, под
# которую считалась вёрстка слайда, и менять её нельзя: другая ширина даёт
# другие переносы. А вот хранить в файле такой большой растр незачем: при
# ширине страницы 13,3 дюйма 1600 точек дают 120 dpi на экране и остаются
# чёткими при печати на A4.
PDF_IMAGE_W = 1600

# Публичный адрес сайта. Нужен, чтобы относительные ссылки деки («../app/»)
# в раздатке вели на живой сайт, а не на localhost, с которого снимался
# слайд. Значение сверяется с самой декой при сборке — если разойдётся,
# скрипт скажет об этом вслух, а не подставит молча чужой адрес.
PUBLIC_SITE = "https://mandaloriancs.github.io/renthub/"

# Шрифт для невидимого текстового слоя.
#
# Встроенные шрифты PDF кириллицы не знают: слой из них дал бы поиск,
# который не находит «перфоратор». Manrope уже лежит в зависимостях —
# им набрана сама дека, и в файл уходит только использованное
# подмножество, около 60 КБ.
TEXT_FONT = ROOT / "node_modules" / "@expo-google-fonts" / "manrope" / "400Regular" / "Manrope_400Regular.ttf"

# Качество JPEG. 92 при отключённом субсэмплинге (4:4:4) — проверено на
# самой мелкой строке деки: увеличенный вдвое кроп неотличим от PNG.
# Ниже 88 у букв появляется ореол, и это видно на таблицах.
PDF_JPEG_Q = 92


def squeeze(path):
    """Снимок → JPEG для вставки в PDF. Втрое легче, разницы не видно.

    Почему не PNG. PyMuPDF вставляет PNG, разжимая его в RGB: палитра на
    256 цветов, которая вдвое уменьшала сам файл снимка, в PDF не доезжает —
    проверено, страница как была 264 КБ, так и осталась. JPEG он кладёт
    потоком как есть, поэтому вес, который мы видим здесь, — это вес,
    который окажется в файле.

    Почему это безопасно для текста. JPEG портит буквы субсэмплингом
    цветности; при 4:4:4 и качестве 92 портить нечего — дека нарисована
    плоскими заливками, а не фотографиями. Сравнение кропа таблицы при
    двукратном увеличении разницы с PNG не показало.

    Не сработает — отдаём файл как есть: лишний мегабайт дешевле, чем
    сборка, упавшая перед защитой.
    """
    try:
        from PIL import Image
    except ImportError:
        return path.read_bytes()

    try:
        with Image.open(path) as img:
            frame = img.convert("RGB")
            if frame.width > PDF_IMAGE_W:
                height = round(frame.height * PDF_IMAGE_W / frame.width)
                frame = frame.resize((PDF_IMAGE_W, height), Image.LANCZOS)

            buf = io.BytesIO()
            frame.save(buf, format="JPEG", quality=PDF_JPEG_Q, subsampling=0, optimize=True)
            return buf.getvalue()
    except Exception as err:  # noqa: BLE001 — причина не важна, важен запасной путь
        print(f"  ! {path.name}: сжать не вышло ({err}), кладу как есть")
        return path.read_bytes()


def looks_blank(path):
    """Снимок пустой? Тогда это не слайд, а несостоявшаяся отрисовка.

    Браузер иногда отдаёт кадр раньше, чем нарисует страницу: получается
    ровный фон нужного цвета, файл на месте, размер правдоподобный. В деке
    это выглядит как потерянный слайд, и заметить его можно только глазами
    — а смотрят раздатку в первый раз обычно судьи.

    Считаем разброс яркости по сетке точек: у настоящего слайда есть
    заголовок и текст, то есть тёмное на светлом. Порог низкий намеренно —
    задача поймать белый лист, а не оценить композицию.
    """
    try:
        from PIL import Image, ImageStat
    except ImportError:
        return False

    try:
        with Image.open(path) as img:
            stat = ImageStat.Stat(img.convert("L").resize((160, 90)))
            return stat.stddev[0] < 4
    except Exception:  # noqa: BLE001 — проверка не должна ронять сборку
        return False


def build_pdf(shots):
    try:
        import fitz  # PyMuPDF
    except ImportError:
        print("\n✗ Нет PyMuPDF. Поставьте: pip install pymupdf\n")
        return False

    doc = fitz.open()

    # Снимок делался в кадре 1920 точек, страница — 960 пунктов шириной.
    # Один коэффициент переводит одно в другое: и картинку, и координаты
    # ссылок поверх неё.
    k = PAGE_W / SHOT_W
    anchors = slide_anchors()

    for _, path, _links, _lines, _vp in shots:
        page = doc.new_page(width=PAGE_W, height=PAGE_H)
        page.insert_image(fitz.Rect(0, 0, PAGE_W, PAGE_H), stream=squeeze(path))

    # Ссылки — вторым проходом, когда все страницы уже созданы.
    #
    # Порядок здесь не стилистический. Переход внутри файла PyMuPDF молча
    # выбрасывает, если страницы-цели ещё нет: на первом слайде кнопка «Что
    # уже работает» ведёт на пятнадцатую страницу, а её в тот момент не
    # существует. Ошибки нет, ссылки тоже — поймано сверкой готового файла,
    # а не выводом сборки.
    font = "manrope" if TEXT_FONT.exists() else None
    if not font:
        print(f"  ! нет {TEXT_FONT.name} — в PDF не будет поиска по тексту")

    for index, (_, _path, links, lines, viewport) in enumerate(shots):
        page = doc[index]

        # Страница мерила себя в своём вьюпорте, снимок сделан в кадре
        # побольше. Слайд центрирован — значит разница легла поровну.
        dx = (SHOT_W - viewport[0]) / 2 if viewport else 0
        dy = (SHOT_H - viewport[1]) / 2 if viewport else 0

        # Невидимый текстовый слой.
        #
        # Страница остаётся картинкой — рисуется она снимком. Поверх
        # ложится тот же текст режимом «не рисовать» (render_mode 3):
        # его не видно, но он ищется, выделяется и копируется, а
        # скринридер его читает. Раздатка перестаёт быть немой.
        #
        # Строки берутся из деки как есть, поэтому слой всегда совпадает
        # с картинкой: оба сняты с одной страницы в один заход.
        for line in lines if font else []:
            text = line.get("s") or ""
            if not text.strip():
                continue

            size = max(4.0, line["f"] * k)
            # Базовая линия примерно на четверти высоты снизу — точнее без
            # метрик шрифта не выйдет, а для выделения этого хватает.
            point = fitz.Point((line["x"] + dx) * k, (line["y"] + dy + line["t"] * 0.78) * k)

            try:
                page.insert_text(
                    point,
                    text,
                    fontname=font,
                    fontfile=str(TEXT_FONT),
                    fontsize=size,
                    render_mode=3,
                )
            except Exception:  # noqa: BLE001 — строка со странным символом не повод падать
                continue

        for link in links:
            box = fitz.Rect(
                (link["x"] + dx) * k,
                (link["y"] + dy) * k,
                (link["x"] + dx + link["w"]) * k,
                (link["y"] + dy + link["t"]) * k,
            )

            href = (link.get("h") or "").strip()
            if not href:
                continue

            if href.startswith("#"):
                # Внутренний якорь — переход по страницам самого файла.
                target = anchors.get(href[1:])
                if target:
                    # to и zoom обязательны, хотя выглядят необязательными:
                    # без точки назначения PyMuPDF 1.28 переход молча не
                    # создаёт — ни ошибки, ни ссылки. Полчаса ушло на поиск
                    # виновного, потому что искали его в номере страницы.
                    page.insert_link({
                        "kind": fitz.LINK_GOTO,
                        "from": box,
                        "page": target - 1,
                        "to": fitz.Point(0, 0),
                        "zoom": 0,
                    })
                continue

            if href.startswith("http://") or href.startswith("https://"):
                page.insert_link({"kind": fitz.LINK_URI, "from": box, "uri": href})
                continue

            # Относительный путь: в раздатке он должен вести на живой сайт.
            # «../app/» рядом с pitch.html — это /renthub/app/ на Pages.
            page.insert_link({
                "kind": fitz.LINK_URI,
                "from": box,
                "uri": PUBLIC_SITE + href.lstrip("./").lstrip("/"),
            })

    # Метаданные видит не только программа: имя файла в мессенджере может
    # смениться на «document (3).pdf», а заголовок останется. Судья, у
    # которого открыто восемь вкладок, находит нашу по названию.
    doc.set_metadata({
        "title": "RentHUB — питч",
        "author": "RentHUB, Кокшетау",
        "subject": "Аренда строительного инструмента между соседями. Питч для защиты.",
        "keywords": "RentHUB, аренда инструмента, P2P, Кокшетау, MVP",
        "creator": "scripts/pitch-slides.py",
    })

    # Оглавление: один уровень, слайд = закладка.
    titles = slide_titles()
    if len(titles) == len(shots):
        doc.set_toc([[level, name, i + 1] for i, (level, name) in enumerate(titles)])
    else:
        # Разошлось — значит разметка деки изменилась, а разбор нет.
        # Молча отдать PDF без оглавления хуже: пропажу заметят на защите.
        print(f"  ! заголовков {len(titles)}, слайдов {len(shots)} — оглавление пропущено")

    # Как файл откроется у судьи.
    #
    # По умолчанию просмотрщик показывает раздатку лентой с прокруткой и
    # прячет оглавление — то есть ровно так, как её листать неудобно.
    # Четыре ключа в каталоге это меняют:
    #
    #   PageLayout /SinglePage — одна страница на экран, как слайды;
    #   PageMode /UseOutlines  — оглавление сразу открыто сбоку;
    #   OpenAction … /Fit      — первый слайд вписан в окно целиком;
    #   Lang (ru-RU)           — язык документа, его читает скринридер.
    #
    # Ключи необязательные: просмотрщик вправе их игнорировать, и ни один
    # из них ничего не ломает, если проигнорирован.
    try:
        catalog = doc.pdf_catalog()
        doc.xref_set_key(catalog, "PageLayout", "/SinglePage")
        doc.xref_set_key(catalog, "PageMode", "/UseOutlines")
        doc.xref_set_key(catalog, "Lang", "(ru-RU)")
        doc.xref_set_key(catalog, "OpenAction", f"<</Type/Action/S/GoTo/D[{doc[0].xref} 0 R /Fit]>>")
    except Exception as err:  # noqa: BLE001 — настройки показа не стоят упавшей сборки
        print(f"  ! не удалось задать вид при открытии ({err})")

    doc.save(str(OUT_PDF), deflate=True, garbage=3)
    doc.close()

    if not numbers_survived():
        return False

    if not qr_works():
        return False

    # Отпечаток исходника рядом с файлом: по нему npm run check:pitch
    # понимает, не отстала ли раздатка от деки. Сравнение по времени
    # файлов не работает — git времени не хранит.
    Path(str(OUT_PDF) + ".sha").write_text(f"{deck_fingerprint()}\n", encoding="utf-8")
    return True


def numbers_survived():
    """Главные цифры деки должны найтись в готовом файле.

    Проверка появилась после дорогой ошибки. Числа в полосе экономики
    «набегают» анимацией за 900 мс; в headless кадры выдаёт композитор,
    которого при съёмке может не быть, и анимация застревает на середине.
    В раздатку уехало «436 ₸» вместо «2 100 ₸» — цифра, которой нет ни в
    разметке, ни в расчётах.

    Ни одна прежняя проверка этого не ловила: npm run check:pitch сверяет
    разметку, а разметка была верной. Врал снимок. Теперь сверяется то,
    что реально попало в файл, — благо текстовый слой это позволяет.

    Сверяются значения полосы (.stat-v): это выводы слайдов о деньгах, и
    именно их трогает анимация.
    """
    try:
        import fitz
    except ImportError:
        return True

    html = DECK.read_text(encoding="utf-8")
    wanted = re.findall(r'<div class="stat-v">([^<]+)</div>', html)
    if not wanted:
        return True

    doc = fitz.open(str(OUT_PDF))
    text = " ".join(doc[i].get_text() for i in range(doc.page_count))
    doc.close()

    # Пробел внутри числа бывает неразрывным, а знак тенге в слое может и
    # не встретиться: сверяем сами цифры.
    only_digits = lambda value: re.sub(r"[^0-9]", "", value)
    haystack = only_digits(text)

    missing = [v.strip() for v in wanted if only_digits(v) and only_digits(v) not in haystack]

    if missing:
        print("\n✗ В готовом PDF нет чисел деки: " + ", ".join(missing))
        print("  Похоже, в кадр попала анимация чисел. Запустите сборку ещё раз.\n")
        return False

    return True


def qr_works():
    """QR со слайда должен читаться камерой, а не просто выглядеть кодом.

    Код рисуется в разметку скриптом npm run qr, потом снимается вместе со
    слайдом, потом жмётся в JPEG и уменьшается до 1600 точек. На любом из
    этих шагов он может стать нечитаемым — и снаружи это не отличить:
    квадрат с узором выглядит одинаково рабочим и сломанным.

    Поэтому проверяем не наличие картинки, а результат: рендерим страницу
    и декодируем то, что получилось, как это сделает камера судьи. Адрес
    должен совпасть с тем, что записано в разметке.

    Нет OpenCV — проверка молча пропускается: она полезная, но не повод
    ронять сборку на машине без лишней библиотеки.
    """
    try:
        import cv2
        import fitz
        import numpy as np
    except ImportError:
        return True

    html = DECK.read_text(encoding="utf-8")
    wanted = re.findall(r'data-qr="([^"]+)"', html)
    if not wanted:
        return True

    doc = fitz.open(str(OUT_PDF))
    detector = cv2.QRCodeDetector()
    seen = set()

    for index in range(doc.page_count):
        # Тройное увеличение: камера смотрит на проектор с расстояния, а
        # детектору нужны различимые ячейки. Меньший масштаб даёт ложные
        # «не прочитал» на странице, которая человеком читается нормально.
        pix = doc[index].get_pixmap(matrix=fitz.Matrix(3, 3))
        frame = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, pix.n)
        frame = cv2.cvtColor(frame, cv2.COLOR_RGBA2BGR if pix.n == 4 else cv2.COLOR_RGB2BGR)

        value, _points, _straight = detector.detectAndDecode(frame)
        if value:
            seen.add(value)

    doc.close()

    missing = [url for url in wanted if url not in seen]

    if missing:
        print("\n✗ QR в PDF не читается или ведёт не туда: " + ", ".join(missing))
        if seen:
            print("  Прочитано вместо этого: " + ", ".join(sorted(seen)))
        print("  Перерисовать код: npm run qr, затем собрать заново.\n")
        return False

    print(f"  QR проверен камерой: {', '.join(sorted(seen))}")
    return True


def build_pptx(shots):
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        print("\n✗ Нет python-pptx. Поставьте: pip install python-pptx\n")
        return False

    deck = Presentation()
    deck.slide_width = Inches(13.333)
    deck.slide_height = Inches(7.5)
    blank = deck.slide_layouts[6]

    for _, path, _links, _lines, _vp in shots:
        slide = deck.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(path), 0, 0, width=deck.slide_width, height=deck.slide_height
        )

    deck.save(str(OUT_PPTX))

    # Отпечаток рядом с файлом — тот же приём, что у PDF: иначе
    # презентация устаревает молча, а замечают это на защите.
    Path(str(OUT_PPTX) + ".sha").write_text(f"{deck_fingerprint()}\n", encoding="utf-8")
    return True


def main(argv):
    want_pdf = "--pdf" in argv or not ("--pdf" in argv or "--pptx" in argv)
    want_pptx = "--pptx" in argv or not ("--pdf" in argv or "--pptx" in argv)

    browser = find_browser()
    if not browser:
        print("\n✗ Не нашёл Chrome или Edge — снимать слайды нечем.")
        for p in BROWSERS:
            print(f"    {p}")
        return 1

    total = len(re.findall(r"<section", DECK.read_text(encoding="utf-8")))
    if total < 5:
        print(f"\n✗ В деке найдено {total} секций — образец сломан.\n")
        return 1

    server = serve()
    try:
        with tempfile.TemporaryDirectory() as tmp:
            shots = []
            for n in range(1, total + 1):
                shot, links, lines, viewport = shoot(browser, tmp, n)
                if not shot:
                    print(f"\n✗ Слайд {n} не снялся.\n")
                    return 1
                if looks_blank(shot):
                    print(f"\n✗ Слайд {n} снялся пустым — браузер не успел его нарисовать.")
                    print("  Запустите сборку ещё раз.\n")
                    return 1

                shots.append((n, shot, links, lines, viewport))
                note = f", ссылок {len(links)}" if links else ""
                print(f"  снят слайд {n:2d} — {round(shot.stat().st_size / 1024)} КБ{note}")

            if want_pdf and not build_pdf(shots):
                return 1
            if want_pptx and not build_pptx(shots):
                return 1
    finally:
        server.shutdown()

    print()
    if want_pdf:
        print(f"✓ Раздатка: landing/{OUT_PDF.name} — {total} страниц, "
              f"{round(OUT_PDF.stat().st_size / 1024)} КБ")
    if want_pptx:
        print(f"✓ Презентация: landing/{OUT_PPTX.name} — {total} слайдов, "
              f"{round(OUT_PPTX.stat().st_size / 1024)} КБ")

    print("\n  Оба файла — снимки landing/pitch.html: текст в них не правится.")
    print("  Числа сверяет npm run check:pitch по разметке той же страницы.\n")
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv[1:]))
