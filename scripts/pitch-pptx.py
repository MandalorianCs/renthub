#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Дека презентацией PowerPoint.

    npm run pitch:pptx

Зачем, если есть PDF. PDF — раздатка: его открывают и читают. PPTX
показывают с проектора, листают стрелками и правят перед защитой, если
организатор попросит. Регламент Терриконовой долины («Шаблон ИНК»)
ожидает именно презентацию, а тайминг в нём расписан по слайдам — то
есть слайд там единица, а не страница.

Как собирается. Из той же landing/pitch.html: каждая секция снимается
картинкой в режиме `?slide=N` и кладётся на слайд 16:9. Второй копии
презентации в проекте нет и не будет — числа деки сверяет
npm run check:pitch, и сверяет он разметку этой страницы.

Отсюда честное ограничение: слайды получаются картинками, текст в них не
редактируется. Это осознанный размен. Собрать PPTX «настоящими» текстовыми
блоками значило бы переписать вёрстку деки средствами PowerPoint —
получилась бы вторая презентация, которая разойдётся с первой в первый же
вечер правок. Здесь же расходиться нечему: картинка снята с той страницы,
которую проверяют.
"""

import io
import re
import subprocess
import sys
import tempfile
import time
from http.server import HTTPServer, SimpleHTTPRequestHandler
from pathlib import Path
from threading import Thread

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")

ROOT = Path(__file__).resolve().parent.parent
LANDING = ROOT / "landing"
OUT = LANDING / "RentHUB-pitch.pptx"
PORT = 8901

# 16:9 в точках PowerPoint: 13.333 × 7.5 дюйма. Снимаем вдвое крупнее,
# чтобы на проекторе не рассыпалось.
SHOT_W, SHOT_H = 1920, 1080

BROWSERS = [
    r"C:/Program Files (x86)/Microsoft/Edge/Application/msedge.exe",
    r"C:/Program Files/Microsoft/Edge/Application/msedge.exe",
    r"C:/Program Files/Google/Chrome/Application/chrome.exe",
    "/usr/bin/google-chrome",
]


def find_browser():
    for path in BROWSERS:
        if Path(path).exists():
            return path
    return None


def serve():
    """Тот же локальный сервер, что у печати PDF, и по той же причине.

    Из file:// Chromium не берёт часть шрифтов и SVG — картинка выходит
    в системном шрифте, и заметно это только в готовом файле.
    """

    class Handler(SimpleHTTPRequestHandler):
        def __init__(self, *args, **kwargs):
            super().__init__(*args, directory=str(LANDING), **kwargs)

        def log_message(self, *args):
            pass

    server = HTTPServer(("127.0.0.1", PORT), Handler)
    Thread(target=server.serve_forever, daemon=True).start()
    return server


def main():
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        print("\n✗ Нет python-pptx. Поставьте: pip install python-pptx\n")
        return 1

    browser = find_browser()
    if not browser:
        print("\n✗ Не нашёл Chrome или Edge — снимать слайды нечем.\n")
        return 1

    deck = (LANDING / "pitch.html").read_text(encoding="utf-8")
    total = len(re.findall(r"<section", deck))

    if total < 5:
        print(f"\n✗ В деке найдено {total} секций — образец сломан.\n")
        return 1

    server = serve()
    shots = []

    try:
        with tempfile.TemporaryDirectory() as tmp:
            for n in range(1, total + 1):
                shot = Path(tmp) / f"slide-{n:02d}.png"

                subprocess.run(
                    [
                        browser,
                        "--headless=new",
                        "--disable-gpu",
                        # Свой профиль обязателен: без него запуск не создаёт
                        # процесс, а передаёт команду открытому браузеру — и
                        # снимок не появляется вовсе.
                        f"--user-data-dir={Path(tempfile.gettempdir()) / 'renthub-pptx'}",
                        "--hide-scrollbars",
                        f"--window-size={SHOT_W},{SHOT_H}",
                        f"--screenshot={shot}",
                        "--virtual-time-budget=6000",
                        f"http://127.0.0.1:{PORT}/pitch.html?slide={n}",
                    ],
                    capture_output=True,
                    timeout=90,
                )

                # Chromium иногда возвращает управление раньше, чем допишет
                # файл: тот же случай, что с печатью PDF.
                for _ in range(20):
                    if shot.exists() and shot.stat().st_size > 0:
                        break
                    time.sleep(0.3)

                if not shot.exists():
                    print(f"\n✗ Слайд {n} не снялся.\n")
                    return 1

                shots.append((n, shot.read_bytes()))
                print(f"  снят слайд {n:2d} — {round(shot.stat().st_size / 1024)} КБ")

            deck_pptx = Presentation()
            deck_pptx.slide_width = Inches(13.333)
            deck_pptx.slide_height = Inches(7.5)
            blank = deck_pptx.slide_layouts[6]

            for n, data in shots:
                slide = deck_pptx.slides.add_slide(blank)
                slide.shapes.add_picture(
                    io.BytesIO(data),
                    0,
                    0,
                    width=deck_pptx.slide_width,
                    height=deck_pptx.slide_height,
                )

            deck_pptx.save(OUT)
    finally:
        server.shutdown()

    size = round(OUT.stat().st_size / 1024)
    print(f"\n✓ Презентация собрана: landing/{OUT.name}")
    print(f"  {len(shots)} слайдов, {size} КБ\n")
    print("  Слайды — картинки с той же страницы, что и PDF: текст в них не")
    print("  правится. Числа сверяет npm run check:pitch по landing/pitch.html.\n")
    return 0


if __name__ == "__main__":
    sys.exit(main())
